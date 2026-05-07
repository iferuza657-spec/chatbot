import os
import re
import requests
from groq import Groq
from dotenv import load_dotenv
from telegram import Update, BotCommand
from telegram.ext import (
    ApplicationBuilder,
    MessageHandler,
    CommandHandler,
    filters,
    ContextTypes
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DocPt

load_dotenv()

TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

client = Groq(api_key=GROQ_API_KEY)

user_histories = {}

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Salom! Men AI yordamchiman 🤖\n\n"
        "📝 Matn yozing — AI javob beradi\n"
        "🎤 Ovozli xabar — AI tinglab javob beradi\n"
        "🖼️ Rasm yuboring — AI tasvirlab beradi\n"
        "📊 Prezentatsiya — 'prezentatsiya yasa: mavzu' yozing\n"
        "📄 Word hujjat — 'word yas: mavzu' yozing\n\n"
        "/help — yordam\n"
        "/reset — suhbatni tozalash"
    )

async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "📌 Buyruqlar:\n\n"
        "/start — Botni boshlash\n"
        "/reset — Suhbat tarixini tozalash\n"
        "/help — Yordam\n\n"
        "📊 Prezentatsiya yasash:\n"
        "  'prezentatsiya yasa: sun quyosh haqida'\n\n"
        "📄 Word hujjat yasash:\n"
        "  'word yasa: rezyume'\n"
    )

async def reset(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    user_histories[user_id] = []
    await update.message.reply_text("Suhbat tarixi tozalandi ✅")

def create_pptx(topic, slides_data):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Rang palitrasi
    DARK_BG = RGBColor(18, 18, 40)
    ACCENT = RGBColor(99, 102, 241)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(200, 200, 220)

    for i, slide_info in enumerate(slides_data):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Fon
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        # Chiziq bezak
        from pptx.util import Emu
        line = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(0.08)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

        if i == 0:
            # Birinchi slayd — sarlavha
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_info.get("title", topic)
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = WHITE

            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
            tf2 = sub_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = slide_info.get("content", "")
            p2.font.size = Pt(22)
            p2.font.color.rgb = ACCENT
        else:
            # Oddiy slayd
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(11), Inches(1))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_info.get("title", "")
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = ACCENT

            content = slide_info.get("content", "")
            lines = content.split("\n")
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
            tf2 = content_box.text_frame
            tf2.word_wrap = True

            for j, line in enumerate(lines):
                if line.strip():
                    p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                    p2.text = f"• {line.strip()}" if not line.startswith("•") else line
                    p2.font.size = Pt(20)
                    p2.font.color.rgb = LIGHT_GRAY
                    p2.space_after = Pt(8)

        # Sahifa raqami
        num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3))
        tf_n = num_box.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.text = str(i + 1)
        p_n.font.size = Pt(12)
        p_n.font.color.rgb = LIGHT_GRAY

    filename = "prezentatsiya.pptx"
    prs.save(filename)
    return filename

def create_docx(topic, content):
    doc = Document()

    # Sarlavha
    title = doc.add_heading(topic, 0)
    title.runs[0].font.color.rgb = RGBColor(99, 102, 241)

    doc.add_paragraph("")

    lines = content.split("\n")
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("##"):
            doc.add_heading(line.replace("##", "").strip(), level=2)
        elif line.startswith("#"):
            doc.add_heading(line.replace("#", "").strip(), level=1)
        elif line.startswith("•") or line.startswith("-") or line.startswith("*"):
            p = doc.add_paragraph(line.lstrip("•-* "), style="List Bullet")
        else:
            doc.add_paragraph(line)

    filename = "hujjat.docx"
    doc.save(filename)
    return filename

async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    user_text = update.message.text.strip()

    # Prezentatsiya so'rovi
    pptx_triggers = ["prezentatsiya yasa", "prezentatsiya qil", "slayd yasa", "pptx yasa"]
    docx_triggers = ["word yasa", "word qil", "hujjat yasa", "docx yasa", "rezyume yasa"]

    if any(user_text.lower().startswith(t) for t in pptx_triggers):
        topic = re.split(r"[:—\-]", user_text, 1)[-1].strip() or user_text
        await update.message.reply_text(f"📊 '{topic}' bo'yicha prezentatsiya tayyorlanmoqda...")
        await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="upload_document")

        prompt = f"""'{topic}' mavzusida 6 ta slayd uchun kontent yarat.
Faqat JSON formatida javob ber, boshqa hech narsa yozma:
[
  {{"title": "Sarlavha slayd", "content": "Qisqa tavsif"}},
  {{"title": "Slayd 2 nomi", "content": "- nuqta 1\\n- nuqta 2\\n- nuqta 3"}},
  ...
]"""

        try:
            resp = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}]
            )
            raw = resp.choices[0].message.content
            raw = re.sub(r"```json|```", "", raw).strip()
            import json
            slides_data = json.loads(raw)
            filename = create_pptx(topic, slides_data)
            with open(filename, "rb") as f:
                await update.message.reply_document(f, filename=filename, caption=f"📊 '{topic}' prezentatsiyasi tayyor!")
        except Exception as e:
            await update.message.reply_text(f"Xatolik: {str(e)}")
        return

    if any(user_text.lower().startswith(t) for t in docx_triggers):
        topic = re.split(r"[:—\-]", user_text, 1)[-1].strip() or user_text
        await update.message.reply_text(f"📄 '{topic}' bo'yicha hujjat tayyorlanmoqda...")
        await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="upload_document")

        prompt = f"""'{topic}' mavzusida to'liq va professional hujjat yoz.
Sarlavhalar uchun # va ## ishlat. Ro'yxatlar uchun - ishlat. O'zbek tilida yoz."""

        try:
            resp = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}]
            )
            content = resp.choices[0].message.content
            filename = create_docx(topic, content)
            with open(filename, "rb") as f:
                await update.message.reply_document(f, filename=filename, caption=f"📄 '{topic}' hujjati tayyor!")
        except Exception as e:
            await update.message.reply_text(f"Xatolik: {str(e)}")
        return

    # Oddiy xabar
    if user_id not in user_histories:
        user_histories[user_id] = [
            {
                "role": "system",
                "content": (
                    "Sen foydali yordamchi botsan. O'zbek tilida javob ber. "
                    "Agar foydalanuvchi biror ilova, dastur yoki vosita haqida so'rasa, "
                    "uning rasmiy havolasini ham yubor."
                )
            }
        ]

    user_histories[user_id].append({"role": "user", "content": user_text})
    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")

    try:
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=user_histories[user_id]
        )
        reply = response.choices[0].message.content
        user_histories[user_id].append({"role": "assistant", "content": reply})
        await update.message.reply_text(reply)
    except Exception as e:
        await update.message.reply_text(f"Xatolik yuz berdi: {str(e)}")

async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    try:
        voice = update.message.voice
        file = await context.bot.get_file(voice.file_id)
        audio_data = requests.get(file.file_path).content
        with open("voice.ogg", "wb") as f:
            f.write(audio_data)
        with open("voice.ogg", "rb") as f:
            transcription = client.audio.transcriptions.create(
                file=("voice.ogg", f.read()),
                model="whisper-large-v3",
            )
        user_text = transcription.text
        await update.message.reply_text(f"🎤 Siz dedingiz: {user_text}")

        if user_id not in user_histories:
            user_histories[user_id] = [{"role": "system", "content": "Sen foydali yordamchi botsan. O'zbek tilida javob ber."}]
        user_histories[user_id].append({"role": "user", "content": user_text})
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=user_histories[user_id]
        )
        reply = response.choices[0].message.content
        user_histories[user_id].append({"role": "assistant", "content": reply})
        await update.message.reply_text(reply)
    except Exception as e:
        await update.message.reply_text(f"Xatolik yuz berdi: {str(e)}")

async def handle_image(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    try:
        photo = update.message.photo[-1]
        file = await context.bot.get_file(photo.file_id)
        caption = update.message.caption or "Bu rasmda nima bor?"
        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": file.file_path}},
                    {"type": "text", "text": caption}
                ]
            }]
        )
        await update.message.reply_text(response.choices[0].message.content)
    except Exception as e:
        await update.message.reply_text(f"Xatolik yuz berdi: {str(e)}")

async def post_init(app):
    await app.bot.set_my_commands([
        BotCommand("start", "Botni boshlash"),
        BotCommand("reset", "Suhbatni tozalash"),
        BotCommand("help", "Yordam"),
    ])

def main():
    app = ApplicationBuilder().token(TELEGRAM_TOKEN).post_init(post_init).build()
    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("reset", reset))
    app.add_handler(CommandHandler("help", help_command))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))
    app.add_handler(MessageHandler(filters.VOICE, handle_voice))
    app.add_handler(MessageHandler(filters.PHOTO, handle_image))
    print("Bot ishlamoqda... ✅")
    app.run_polling()

if __name__ == "__main__":
    main()